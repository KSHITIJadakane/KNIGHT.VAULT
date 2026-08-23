import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_top_tier_presentation(output_path="Midnight_Privacy_Payment_Vault_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Institutional Palette (KNIGHT.VAULT Signature Palette)
    BG_COLOR = RGBColor(8, 9, 13)          # #08090d Dark Obsidian
    CARD_BG = RGBColor(14, 17, 23)         # #0e1117 Surface Card
    CARD_BG_ALT = RGBColor(18, 22, 30)     # #12161e Elevated Card
    BORDER_COLOR = RGBColor(35, 42, 56)    # #232a38 Border
    
    TEXT_WHITE = RGBColor(255, 255, 255)   # #ffffff Pure White
    TEXT_PLATINUM = RGBColor(241, 245, 249)# #f1f5f9 Crisp Text
    TEXT_MUTED = RGBColor(160, 174, 192)   # #a0aec0 Slate Muted
    TEXT_DIM = RGBColor(110, 126, 150)     # #6e7e96 Subtle Label
    
    MINT = RGBColor(0, 245, 160)           # #00f5a0 Cyber Mint Primary
    INDIGO = RGBColor(129, 140, 248)       # #818cf8 Soft Indigo
    CYAN = RGBColor(56, 189, 248)          # #38bdf8 Electric Cyan
    ROSE = RGBColor(244, 63, 94)           # #f43f5e Warning Rose
    EMERALD = RGBColor(16, 185, 129)       # #10b981 Success

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="MIDNIGHT NETWORK // ZERO-KNOWLEDGE PROTOCOL"):
        # Top Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = MINT

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

        # Divider Line with Mint Accent Pip
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.733), Inches(0.018))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_COLOR
        line.line.fill.background()

        pip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(1.5), Inches(0.03))
        pip.fill.solid()
        pip.fill.fore_color.rgb = MINT
        pip.line.fill.background()

    def add_card(slide, left, top, width, height, title, items, tag=None, accent_color=MINT, bg_color=CARD_BG, font_size=12, item_spacing=8):
        # Card Surface Container
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)

        # Top Accent Tab
        accent_tab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.3), Inches(top), Inches(min(width - 0.6, 2.2)), Inches(0.045))
        accent_tab.fill.solid()
        accent_tab.fill.fore_color.rgb = accent_color
        accent_tab.line.fill.background()

        # Tag / Subtitle if provided
        top_offset = 0.22
        if tag:
            tb_tag = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + top_offset), Inches(width - 0.6), Inches(0.3))
            tf_tag = tb_tag.text_frame
            p_t = tf_tag.paragraphs[0]
            p_t.text = tag.upper()
            p_t.font.name = "Arial"
            p_t.font.size = Pt(9)
            p_t.font.bold = True
            p_t.font.color.rgb = accent_color
            top_offset += 0.3

        # Title
        tb_title = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + top_offset), Inches(width - 0.6), Inches(0.5))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p = tf_title.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        top_offset += 0.52

        # Bullets / Content
        tb_items = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + top_offset), Inches(width - 0.6), Inches(height - top_offset - 0.15))
        tf_items = tb_items.text_frame
        tf_items.word_wrap = True
        
        for i, item in enumerate(items):
            p_item = tf_items.add_paragraph() if i > 0 else tf_items.paragraphs[0]
            p_item.text = "•  " + item
            p_item.font.name = "Arial"
            p_item.font.size = Pt(font_size)
            p_item.font.color.rgb = TEXT_MUTED
            p_item.space_after = Pt(item_spacing)

    # =========================================================================
    # SLIDE 1: Hero Cover Slide (Institutional Keynote)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    # Top Tag
    tb_tag = s1.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.733), Inches(0.35))
    tf_tag = tb_tag.text_frame
    p0 = tf_tag.paragraphs[0]
    p0.text = "MIDNIGHT NETWORK // ZERO-KNOWLEDGE SETTLEMENT PROTOCOL"
    p0.font.name = "Arial"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = MINT

    # Title & Subtitle
    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "KNIGHT.VAULT"
    p1.font.name = "Arial"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(8)

    p2 = tf1.add_paragraph()
    p2.text = "Institutional Zero-Knowledge Custody, Compact Smart Contracts, ProofStation 0-Gas Balancing, and Armored State Privacy."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_MUTED

    # 4-Column Stat Cards (Fills full 11.733 inches)
    stat_data = [
        ("TARGET CHAIN", "Midnight Preprod", "Live Partnerchain Network", MINT),
        ("CIRCUIT ENGINE", "Compact Circuits", "Halo2 SNARK Synthesizer", INDIGO),
        ("PROOF SERVER", "Docker :6300", "100% Client-Side Proving", CYAN),
        ("USER GAS COST", "0.00 NIGHT", "Sponsored Dust Relay Flow", EMERALD)
    ]
    for i, (label, val, desc, col) in enumerate(stat_data):
        card_x = 0.8 + i * 3.0
        c = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(3.7), Inches(2.78), Inches(1.75))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = BORDER_COLOR
        c.line.width = Pt(1)

        tb_s = s1.shapes.add_textbox(Inches(card_x + 0.2), Inches(3.85), Inches(2.38), Inches(1.4))
        tfs = tb_s.text_frame
        p_l = tfs.paragraphs[0]
        p_l.text = label
        p_l.font.name = "Arial"
        p_l.font.size = Pt(9)
        p_l.font.bold = True
        p_l.font.color.rgb = TEXT_DIM
        p_l.space_after = Pt(4)

        p_v = tfs.add_paragraph()
        p_v.text = val
        p_v.font.name = "Arial"
        p_v.font.size = Pt(14)
        p_v.font.bold = True
        p_v.font.color.rgb = col
        p_v.space_after = Pt(3)

        p_d = tfs.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

    # Architect Signature Bottom Banner
    sig_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.733), Inches(1.3))
    sig_card.fill.solid()
    sig_card.fill.fore_color.rgb = CARD_BG_ALT
    sig_card.line.color.rgb = BORDER_COLOR
    sig_card.line.width = Pt(1)

    tb_sig = s1.shapes.add_textbox(Inches(1.1), Inches(5.75), Inches(11.2), Inches(1.05))
    tf_sig = tb_sig.text_frame
    p_s1 = tf_sig.paragraphs[0]
    p_s1.text = "Architected & Engineered by: Kshitij Adakane (Vibe Coder, Systems Builder & Applied AI/ML)"
    p_s1.font.name = "Arial"
    p_s1.font.size = Pt(12)
    p_s1.font.bold = True
    p_s1.font.color.rgb = MINT
    p_s1.space_after = Pt(4)

    p_s2 = tf_sig.add_paragraph()
    p_s2.text = "Stack: Compact Smart Contracts • Halo2 Proofs • Midnight.js SDK • 1AM Wallet • React 18 • Vite • Tailwind Kinetic UI"
    p_s2.font.name = "Arial"
    p_s2.font.size = Pt(10.5)
    p_s2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: The Public Blockchain Privacy Trap vs KNIGHT.VAULT
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "The Trillion-Dollar Web3 Privacy Gap", "01 / MARKET CONTEXT & PROBLEM")

    add_card(s2, 0.8, 1.75, 5.75, 5.2, "⚠️ Transparent Chains (ETH, BTC, Solana)", [
        "Total Financial Exposure: All balances, employee payrolls, treasury movements, and vendor invoices are publicly browsable forever.",
        "Gas-Trail De-Anonymization: Purchasing native gas tokens creates an irreversible KYC link to every on-chain move.",
        "Regulatory Backlash on Mixers: Legacy mixing pools lack selective compliance disclosures and face universal government bans.",
        "Enterprise Blockade: Global enterprises cannot execute confidential settlements on public blockchains without leaking corporate data."
    ], tag="THE TRANSPARENCY TRAP", accent_color=ROSE, font_size=12, item_spacing=12)

    add_card(s2, 6.78, 1.75, 5.75, 5.2, "🛡️ KNIGHT.VAULT Solution", [
        "Witness Isolation: Owner secret keys and withdrawal authorizations remain isolated in local client memory. Only succinct ZK proofs touch chain state.",
        "Sponsored Dust Relay: Payers check out seamlessly without acquiring testnet gas tokens or solving faucet captchas.",
        "Deterministic Ledger Accounting: Public counters track verified settlement totals while protecting transactor identities.",
        "Selective Verifiable Receipts: QR cryptographic receipts allow instant auditability without deanonymizing balances."
    ], tag="THE ZERO-KNOWLEDGE SHIELD", accent_color=MINT, font_size=12, item_spacing=12)

    # =========================================================================
    # SLIDE 3: Compact Smart Contract Architecture
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Compact Smart Contract Architecture", "02 / ZERO-KNOWLEDGE CIRCUITS")

    # Left Code Box
    c_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(6.8), Inches(5.2))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = CARD_BG
    c_box.line.color.rgb = BORDER_COLOR
    c_box.line.width = Pt(1)

    tb_code = s3.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(6.4), Inches(4.85))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    
    code_lines = [
        ("// payment.compact (Midnight Network stdlib)", TEXT_DIM),
        ("contract PaymentVault {", TEXT_WHITE),
        ("  ledger balance: Counter;", MINT),
        ("  ledger totalDeposited: Counter;", MINT),
        ("  ledger ownerCommitment: Bytes<32>;", MINT),
        ("", TEXT_WHITE),
        ("  // Public Payer Inflow (Zero User Gas)", TEXT_DIM),
        ("  export circuit receiveUnshielded(amount: Uint<64>): Void {", CYAN),
        ("    balance.increment(amount);", TEXT_PLATINUM),
        ("    totalDeposited.increment(amount);", TEXT_PLATINUM),
        ("  }", CYAN),
        ("", TEXT_WHITE),
        ("  // Confidential Withdrawal (ZK Owner Auth)", TEXT_DIM),
        ("  export circuit withdraw(amount: Uint<64>, recipient: ...): Void {", INDIGO),
        ("    witness ownerSecretKey: Bytes<32>;", MINT),
        ("    assert persistentCommit(ownerSecretKey) == ownerCommitment;", MINT),
        ("    balance.decrement(amount);", TEXT_PLATINUM),
        ("    sendUnshielded(amount, recipient);", TEXT_PLATINUM),
        ("  }", INDIGO),
        ("}", TEXT_WHITE)
    ]
    for i, (line_text, col) in enumerate(code_lines):
        p_c = tf_code.add_paragraph() if i > 0 else tf_code.paragraphs[0]
        p_c.text = line_text
        p_c.font.name = "Consolas"
        p_c.font.size = Pt(9.5)
        p_c.font.color.rgb = col
        p_c.space_after = Pt(2)

    # Right 3 Architecture Cards
    add_card(s3, 7.8, 1.75, 4.733, 1.62, "1. receiveUnshielded()", [
        "Inflow entrypoint. Increments public ledger balance.",
        "Compatible with sponsored dust transactions for zero-gas checkouts."
    ], tag="INFLOW CIRCUIT", accent_color=CYAN, font_size=10.5, item_spacing=4)

    add_card(s3, 7.8, 3.54, 4.733, 1.62, "2. withdraw()", [
        "Outflow entrypoint. Verifies secret key via private witness proof.",
        "Funds transferred unshielded to designated recipient address."
    ], tag="OUTFLOW CIRCUIT", accent_color=INDIGO, font_size=10.5, item_spacing=4)

    add_card(s3, 7.8, 5.33, 4.733, 1.62, "3. Deterministic Ledger", [
        "On-chain Counter types provide clean state isolation.",
        "Live indexing via Midnight GraphQL v4 websocket subscriptions."
    ], tag="PUBLIC STATE", accent_color=MINT, font_size=10.5, item_spacing=4)

    # =========================================================================
    # SLIDE 4: 4-Stage Runtime Settlement Lifecycle
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "4-Stage Zero-Knowledge Settlement Pipeline", "03 / RUNTIME EXECUTION LIFECYCLE")

    pipeline_cards = [
        ("STAGE 01", "Private Witness", [
            "Client runtime encapsulates secret key and transaction parameters.",
            "Witness data stays in local memory and is never broadcast across network."
        ], MINT),
        ("STAGE 02", "Halo2 Synthesis", [
            "Local Docker ProofStation (:6300) runs proving key synthesis.",
            "Generates succinct Zero-Knowledge proof transcript in seconds."
        ], INDIGO),
        ("STAGE 03", "Relay Balancing", [
            "ProofStation attaches required unshielded UTXO dust collateral.",
            "Eliminates user gas fees completely without requiring faucet tokens."
        ], CYAN),
        ("STAGE 04", "Ledger Finality", [
            "Midnight Partnerchain validators verify proof validity.",
            "GraphQL indexer captures state transition and updates vault liquidity live."
        ], EMERALD)
    ]
    for i, (stage_tag, title, items, col) in enumerate(pipeline_cards):
        add_card(s4, 0.8 + i * 2.98, 1.75, 2.78, 5.2, title, items, tag=stage_tag, accent_color=col, font_size=11.5, item_spacing=10)

    # =========================================================================
    # SLIDE 5: Security Benchmark Matrix
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Security & Privacy Benchmark Matrix", "04 / COMPARATIVE ANALYSIS")

    matrix_rows = [
        ("Owner Identity Visibility", "Publicly linkable to wallet address", "Shielded via private ZK witness key"),
        ("Payer Friction", "Requires native gas token in every wallet", "Zero-gas sponsored checkout flow"),
        ("Auditability & Compliance", "Scrapes block explorer (zero privacy)", "Selective cryptographic QR receipts"),
        ("Smart Contract Engine", "EVM Bytecode (Transparent Execution)", "Compact Circuits (ZK Proof Compilation)"),
        ("Proving Architecture", "N/A (No Client Proving)", "100% Local ProofStation (Docker :6300)")
    ]

    # Table Container
    t_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(11.733), Inches(5.2))
    t_box.fill.solid()
    t_box.fill.fore_color.rgb = CARD_BG
    t_box.line.color.rgb = BORDER_COLOR
    t_box.line.width = Pt(1)

    # Table Headers
    tb_th = s5.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(11.1), Inches(0.4))
    tf_th = tb_th.text_frame
    p_th = tf_th.paragraphs[0]
    p_th.text = f"{'FEATURE / CAPABILITY':<32} {'TRADITIONAL PUBLIC VAULT':<34} {'KNIGHT.VAULT (MIDNIGHT)'}"
    p_th.font.name = "Consolas"
    p_th.font.size = Pt(11)
    p_th.font.bold = True
    p_th.font.color.rgb = MINT

    # Separator
    sep = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(2.45), Inches(11.1), Inches(0.015))
    sep.fill.solid()
    sep.fill.fore_color.rgb = BORDER_COLOR
    sep.line.fill.background()

    # Rows
    for i, (feat, pub_val, zkp_val) in enumerate(matrix_rows):
        row_y = 2.65 + i * 0.8
        r_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(row_y), Inches(11.1), Inches(0.68))
        r_box.fill.solid()
        r_box.fill.fore_color.rgb = CARD_BG_ALT
        r_box.line.color.rgb = BORDER_COLOR
        r_box.line.width = Pt(0.5)

        tb_r = s5.shapes.add_textbox(Inches(1.25), Inches(row_y + 0.14), Inches(10.8), Inches(0.45))
        tf_r = tb_r.text_frame
        p_r = tf_r.paragraphs[0]
        p_r.text = f"• {feat:<28} | {pub_val:<33} | ✓ {zkp_val}"
        p_r.font.name = "Consolas"
        p_r.font.size = Pt(10)
        p_r.font.color.rgb = TEXT_PLATINUM

    # =========================================================================
    # SLIDE 6: Verifiability Matrix // Real vs Bluff
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Verifiability Matrix: Real Execution vs Bluff", "05 / CRYPTOGRAPHIC AUDITABILITY")

    add_card(s6, 0.8, 1.75, 2.75, 5.2, "01. Docker SNARK Prover", [
        "Real ProofStation container (:6300).",
        "Evaluates real mathematical constraints.",
        "Proving key: withdraw.prover (2.82 MB).",
        "Inspect via: docker logs proof-server.",
        "Zero mock transcripts."
    ], tag="PROOF ENGINE", accent_color=MINT, font_size=11, item_spacing=8)

    add_card(s6, 3.78, 1.75, 2.75, 5.2, "02. Contract Hash", [
        "Deterministic 64-char address.",
        "Compiled from payment.compact.",
        "Lace owner secret witness commitment.",
        "dd9c5ee171fbb3516cd2...",
        "Unique on-chain contract state."
    ], tag="ON-CHAIN STATE", accent_color=INDIGO, font_size=11, item_spacing=8)

    add_card(s6, 6.78, 1.75, 2.75, 5.2, "03. Lace Preprod UTXO", [
        "Live wallet funded with 1,000 tNIGHT.",
        "Tx ID: 4b853703926172159c...",
        "Preprod network attachment.",
        "Native stake capacity for tDUST.",
        "Real cryptographic identity."
    ], tag="WALLET INTEGRATION", accent_color=CYAN, font_size=11, item_spacing=8)

    add_card(s6, 9.78, 1.75, 2.75, 5.2, "04. Zero-Gas Relay", [
        "Executed with 0 / 0 tDUST in Lace.",
        "Sponsored relayer bypasses gas tank.",
        "Zero user cost (0.00 NIGHT).",
        "Contract custody holds 35 tNIGHT.",
        "Proof of zero-friction DeFi."
    ], tag="GAS SPONSORSHIP", accent_color=ROSE, font_size=11, item_spacing=8)

    # =========================================================================
    # SLIDE 7: Roadmap & Developer Summary (Kshitij Adakane)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Production Roadmap & Architect Profile", "06 / PROTOCOL FUTURE & TEAM")

    add_card(s7, 0.8, 1.75, 3.75, 3.4, "Phase 1: Preprod (Live)", [
        "Full Lace & 1AM connector integrations.",
        "Compact smart contracts on Preprod.",
        "Sponsored dust balancing engine.",
        "ProofStation client-side proving (:6300)."
    ], tag="LIVE PRODUCTION", accent_color=MINT, font_size=11.5, item_spacing=8)

    add_card(s7, 4.78, 1.75, 3.75, 3.4, "Phase 2: Multi-Sig & Vesting", [
        "m-of-n threshold witness approval circuits.",
        "Time-locked payment releases.",
        "Regulatory viewing keys for tax & compliance.",
        "Native Midnight mobile SDK."
    ], tag="NEXT MILESTONE", accent_color=INDIGO, font_size=11.5, item_spacing=8)

    add_card(s7, 8.78, 1.75, 3.75, 3.4, "Phase 3: Mainnet Ecosystem", [
        "Cardano Partnerchain bridge deployment.",
        "Cross-chain shielded atomic swaps.",
        "Decentralized merchant payment widget.",
        "Zero-knowledge escrow dispute resolution."
    ], tag="ENTERPRISE SCALE", accent_color=CYAN, font_size=11.5, item_spacing=8)

    # Master Architect Card
    arch_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.35), Inches(11.733), Inches(1.6))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = CARD_BG_ALT
    arch_box.line.color.rgb = BORDER_COLOR
    arch_box.line.width = Pt(1)

    tb_a = s7.shapes.add_textbox(Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.3))
    tf_a = tb_a.text_frame
    p_a1 = tf_a.paragraphs[0]
    p_a1.text = "KNIGHT.VAULT Protocol Architect: Kshitij Adakane"
    p_a1.font.name = "Arial"
    p_a1.font.size = Pt(14)
    p_a1.font.bold = True
    p_a1.font.color.rgb = MINT
    p_a1.space_after = Pt(4)

    p_a2 = tf_a.add_paragraph()
    p_a2.text = "Vibe Coder, Systems Builder & Applied AI/ML • Automotive & Embedded Systems Engineering • Applied Mathematics"
    p_a2.font.name = "Arial"
    p_a2.font.size = Pt(11)
    p_a2.font.color.rgb = TEXT_WHITE
    p_a2.space_after = Pt(4)

    p_a3 = tf_a.add_paragraph()
    p_a3.text = "Open Source Repository: https://github.com/KSHITIJadakane/KNIGHT.VAULT • Built for Midnight Network Hackathon"
    p_a3.font.name = "Arial"
    p_a3.font.size = Pt(10)
    p_a3.font.color.rgb = TEXT_MUTED

    prs.save(output_path)
    print(f"[SUCCESS] Top-Tier Pitch Deck saved to: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "Midnight_Privacy_Payment_Vault_Presentation.pptx"
    create_top_tier_presentation(out)
